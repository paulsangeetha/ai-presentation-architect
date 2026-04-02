import streamlit as st
import requests
import json
import io
import re
from pptx import Presentation

# ---------------- CONFIG ----------------
st.set_page_config(page_title="AI Presentation Architect", layout="wide")

CHAT_API_URL = "https://api.sarvam.ai/v1/chat/completions"
TRANSLATE_API_URL = "https://api.sarvam.ai/translate"

SUPPORTED_LANGUAGES = {
    'en-IN': 'English',
    'hi-IN': 'Hindi',
    'ta-IN': 'Tamil',
    'te-IN': 'Telugu',
    'bn-IN': 'Bengali',
    'kn-IN': 'Kannada',
    'mr-IN': 'Marathi',
    'gu-IN': 'Gujarati'
}

# ---------------- HELPERS ----------------

def clean_text(text):
    return re.sub(r"```json|```", "", text).strip()


def safe_json_parse(text):
    try:
        return json.loads(text)
    except Exception:
        return None


# ---------------- AI FUNCTIONS ----------------

def generate_english_presentation(topic, api_key, slide_count):
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }

    prompt = f"""
    Generate a {slide_count}-slide presentation on '{topic}'.
    Return ONLY JSON array like:
    [
      {{
        "title": "Title",
        "content": "• point1\\n• point2\\n• point3"
      }}
    ]
    """

    payload = {
        "model": "sarvam-m",
        "messages": [{"role": "user", "content": prompt}]
    }

    try:
        response = requests.post(CHAT_API_URL, headers=headers, json=payload)

        # 🔍 DEBUG INFO IN UI
        st.write("📡 API Status:", response.status_code)
        st.write("📄 Raw Response:", response.text[:1000])

        if response.status_code != 200:
            raise Exception(f"API Error: {response.text}")

        try:
            data = response.json()
        except Exception:
            raise Exception(f"❌ Not JSON:\n{response.text}")

        # Safe extraction
        content = (
            data.get("choices", [{}])[0]
            .get("message", {})
            .get("content", "")
        )

        if not content:
            raise Exception("❌ Empty content from API")

        content = clean_text(content)

        slides = safe_json_parse(content)

        if slides is None:
            raise Exception(f"❌ JSON parsing failed:\n{content}")

        return slides

    except Exception as e:
        raise Exception(f"Presentation Generation Failed:\n{e}")


def translate_content(text, target_lang, api_key):
    if not text.strip() or target_lang == "en-IN":
        return text

    headers = {
        "api-subscription-key": api_key,
        "Content-Type": "application/json"
    }

    payload = {
        "input": text,
        "source_language_code": "en-IN",
        "target_language_code": target_lang
    }

    try:
        response = requests.post(TRANSLATE_API_URL, headers=headers, json=payload)

        st.write("🌐 Translate Status:", response.status_code)

        if response.status_code != 200:
            return text  # fallback

        data = response.json()
        return data.get("translated_text", text)

    except Exception:
        return text


def create_ppt(slides, topic, lang):
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = topic
    slide.placeholders[1].text = f"Generated in {lang}"

    # Content slides
    for s in slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = s.get("title", "Title")
        slide.placeholders[1].text = s.get("content", "")

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer


# ---------------- UI ----------------

st.title("🤖 AI Presentation Architect (Debug Version)")

if "stage" not in st.session_state:
    st.session_state.stage = "input"


# -------- INPUT --------
if st.session_state.stage == "input":

    st.header("Step 1: Input Details")

    api_key = st.text_input("Enter Sarvam API Key", type="password")
    topic = st.text_input("Enter Topic")

    language = st.selectbox(
        "Select Language",
        list(SUPPORTED_LANGUAGES.keys()),
        format_func=lambda x: SUPPORTED_LANGUAGES[x]
    )

    slide_count = st.slider("Slides", 2, 10, 5)

    if st.button("Generate"):

        if not api_key or not topic:
            st.warning("Please enter API key and topic")
            st.stop()

        with st.spinner("Generating..."):

            try:
                st.info("Generating English slides...")
                english_slides = generate_english_presentation(
                    topic, api_key, slide_count
                )

                st.info("Translating...")
                translated_topic = translate_content(topic, language, api_key)

                translated_slides = []
                for s in english_slides:
                    translated_slides.append({
                        "title": translate_content(s["title"], language, api_key),
                        "content": translate_content(s["content"], language, api_key)
                    })

                st.info("Creating PPT...")
                ppt = create_ppt(
                    translated_slides,
                    translated_topic,
                    SUPPORTED_LANGUAGES[language]
                )

                st.session_state.ppt = ppt
                st.session_state.file = f"{topic.replace(' ', '_')}.pptx"
                st.session_state.stage = "download"
                st.rerun()

            except Exception as e:
                st.error(str(e))


# -------- DOWNLOAD --------
if st.session_state.stage == "download":

    st.header("✅ Download Your Presentation")

    st.download_button(
        "Download PPT",
        data=st.session_state.ppt,
        file_name=st.session_state.file,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    if st.button("Start Again"):
        st.session_state.stage = "input"
        st.rerun()